#!/usr/bin/env python3
"""Generate strategic-reporting-deck.pptx from outline content.

Requires: python-pptx (use project venv: ../.venv-pptx)

  cd docs && ../.venv-pptx/bin/python generate_strategic_reporting_deck.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "strategic-reporting-deck.pptx"

SLIDES = [
    (
        "Strategic reporting & analytics",
        [
            "Nextiva vs. market (Dialpad, RingCentral EX/CX)",
            "Canned reports · Custom builder · API / warehouse / BI",
            "March 2026",
        ],
    ),
    (
        "Executive thesis",
        [
            "Three coordinated capabilities: canned · builder · API/warehouse",
            "One metrics layer: same numbers in UI, export, and API",
            "MultiCaaS → unified EX + CX analytics is table stakes",
        ],
    ),
    (
        "Strategy & approach",
        [
            "Trustworthy metrics + progressive disclosure",
            "Enterprise egress first-class (retention, formats, connectors)",
            "Rich operational analytics in-app; deep modeling in customer BI",
        ],
    ),
    (
        "Target architecture (logical)",
        [
            "Sources → Ingest (events, AI enrich) → Metrics catalog",
            "Experiences: canned · dashboards · report builder",
            "Egress: REST API · bulk export · warehouse connector · Power BI",
            "See strategic-reporting-complete-guide.md — Mermaid §4.1",
        ],
    ),
    (
        "Data flow: platform → customer BI",
        [
            "Curated facts & metric marts inside Nextiva",
            "API / bulk / connector → customer warehouse or lake",
            "dbt / orchestration → executive dashboards",
            "See complete guide — Mermaid §4.2",
        ],
    ),
    (
        "Flows — personas",
        [
            "Supervisor: wallboard + canned shift report + schedule",
            "Analyst: template → metrics → filters → save/share",
            "Data engineer: connector/API → warehouse → BI",
            "See complete guide — Mermaid §5",
        ],
    ),
    (
        "RCA — root cause analysis",
        [
            "Symptoms: RFP friction, PS-heavy custom reports, AI narrative gap",
            "Roots: semantic immaturity, persona gaps, egress secondary, EX/CX silo",
            "Corrective: metrics catalog, builder MVP, AI canned packs, warehouse MVP",
        ],
    ),
    (
        "FBA — features, benefits, attributes",
        [
            "Canned: breadth + AI packs → TTV (↑ activation, adoption)",
            "Builder: templates + catalog → less PS (% from template)",
            "Egress: connector + retention → enterprise readiness (SLAs, audit)",
        ],
    ),
    (
        "Functional blocks",
        [
            "B1 Capture → B2 Process → B3 Metrics",
            "B4 Canned · B5 Builder → B6 Distribution",
            "B7 API · B8 Warehouse → stakeholders & customer BI",
        ],
    ),
    (
        "Phased roadmap",
        [
            "P1 Foundation: canned depth, retention, metrics glossary",
            "P2 Parity: flagship warehouse connector, AI canned, API hardening",
            "P3 Differentiation: visual builder, large metric library, unified EX+CX",
            "P4 Platform: streaming, more connectors, predictive",
        ],
    ),
    (
        "Action items (prioritize)",
        [
            "A1 Metrics catalog v0 + owners",
            "A2 Win/loss: top reporting objections",
            "A3 Canned gap list vs competitors",
            "A4 Retention policy target 90d+",
            "A5 Warehouse connector PRD (e.g. Snowflake)",
        ],
    ),
    (
        "Competitive snapshot",
        [
            "RingCentral: breadth (250+ reports, 350+ metrics class)",
            "Dialpad: operational AI (sentiment, floor insight)",
            "Nextiva: solid core — close gaps on AI canned, builder, connectors",
        ],
    ),
    (
        "Next steps",
        [
            "Assign owners — full table in strategic-reporting-complete-guide.md §10",
            "Export Mermaid diagrams from guide as PNG for visuals",
            "Align SE demo script with metrics glossary + warehouse story",
        ],
    ),
]


def add_title_body(prs: Presentation, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    tslide = prs.slides.add_slide(prs.slide_layouts[0])
    tslide.shapes.title.text = SLIDES[0][0]
    subtitle = tslide.placeholders[1]
    subtitle.text = "\n".join(SLIDES[0][1])

    for title, bullets in SLIDES[1:]:
        add_title_body(prs, title, bullets)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
